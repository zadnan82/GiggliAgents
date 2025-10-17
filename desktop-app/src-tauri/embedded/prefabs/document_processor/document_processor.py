"""
Document Processor Prefab
Extracts text from various file formats
"""

from pathlib import Path
from typing import List
import PyPDF2
from docx import Document
from pptx import Presentation


class DocumentProcessor:
    """Process documents into text chunks"""

    def __init__(self, chunk_size: int = 500, chunk_overlap: int = 50):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    def extract_text(self, file_path: str) -> List[str]:
        """
        Extract text from document and split into chunks

        Args:
            file_path: Path to document

        Returns:
            List of text chunks
        """
        path = Path(file_path)
        extension = path.suffix.lower()

        print(f"📄 Processing: {path.name}")

        # Extract based on file type
        if extension == ".pdf":
            text = self._extract_pdf(path)
        elif extension in [".docx", ".doc"]:
            text = self._extract_docx(path)
        elif extension in [".pptx", ".ppt"]:
            text = self._extract_pptx(path)
        elif extension in [".txt", ".md"]:
            text = self._extract_txt(path)
        elif extension in [".xlsx", ".xls"]:
            text = self._extract_excel(path)
        elif extension == ".csv":
            text = self._extract_csv(path)
        elif extension in [".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".gif"]:
            text = self._extract_image_ocr(path)
        elif extension in [".mp3", ".wav", ".m4a", ".flac", ".ogg"]:
            text = self._extract_audio(path)
        elif extension in [".mp4", ".avi", ".mov", ".mkv"]:
            text = self._extract_video(path)
        elif extension == ".zip":
            text = self._extract_zip(path)
        elif extension == ".7z":
            text = self._extract_7z(path)
        elif extension in [".html", ".htm"]:
            text = self._extract_html(path)
        elif extension == ".json":
            text = self._extract_json(path)
        elif extension == ".xml":
            text = self._extract_xml(path)
        else:
            raise ValueError(f"Unsupported file type: {extension}")

        # Split into chunks
        chunks = self._chunk_text(text)

        print(f"✅ Extracted {len(chunks)} chunks from {path.name}")

        return chunks

    # ============================================
    # EXISTING EXTRACTORS
    # ============================================

    def _extract_pdf(self, path: Path) -> str:
        """Extract text from PDF"""
        text = ""
        with open(path, "rb") as f:
            pdf = PyPDF2.PdfReader(f)
            for page in pdf.pages:
                text += page.extract_text() + "\n\n"
        return text

    def _extract_docx(self, path: Path) -> str:
        """Extract text from DOCX"""
        doc = Document(path)
        text = "\n\n".join([para.text for para in doc.paragraphs])
        return text

    def _extract_pptx(self, path: Path) -> str:
        """Extract text from PPTX"""
        prs = Presentation(path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n\n"
        return text

    def _extract_txt(self, path: Path) -> str:
        """Extract text from TXT/MD"""
        with open(path, "r", encoding="utf-8") as f:
            return f.read()

    # ============================================
    # EXCEL & CSV
    # ============================================

    def _extract_excel(self, path: Path) -> str:
        """Extract text from Excel"""
        try:
            import openpyxl

            workbook = openpyxl.load_workbook(path, data_only=True)
            text = ""

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"\n\n=== Sheet: {sheet_name} ===\n\n"

                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(
                        [str(cell) if cell is not None else "" for cell in row]
                    )
                    if row_text.strip():
                        text += row_text + "\n"

            return text
        except Exception as e:
            print(f"⚠️ Excel extraction failed, trying pandas: {e}")
            # Fallback to pandas
            return self._extract_excel_pandas(path)

    def _extract_excel_pandas(self, path: Path) -> str:
        """Extract Excel using pandas (fallback)"""
        import pandas as pd

        text = ""
        xls = pd.ExcelFile(path)

        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name)
            text += f"\n\n=== Sheet: {sheet_name} ===\n\n"
            text += df.to_string(index=False) + "\n"

        return text

    def _extract_csv(self, path: Path) -> str:
        """Extract text from CSV"""
        import pandas as pd

        try:
            df = pd.read_csv(path)
            return df.to_string(index=False)
        except Exception as e:
            print(f"⚠️ Pandas failed, trying basic CSV: {e}")
            # Fallback to basic CSV
            import csv

            text = ""
            with open(path, "r", encoding="utf-8") as f:
                reader = csv.reader(f)
                for row in reader:
                    text += " | ".join(row) + "\n"
            return text

    # ============================================
    # IMAGES (OCR)
    # ============================================

    def _extract_image_ocr(self, path: Path) -> str:
        """Extract text from image using OCR"""
        try:
            from PIL import Image
            import pytesseract

            print(f"🖼️ Running OCR on {path.name}...")

            image = Image.open(path)
            text = pytesseract.image_to_string(image)

            if not text.strip():
                return f"[Image: {path.name} - No text detected]"

            return f"=== OCR from {path.name} ===\n\n{text}"

        except Exception as e:
            print(f"❌ OCR failed: {e}")
            return f"[Image: {path.name} - OCR failed: {e}]"

    # ============================================
    # AUDIO & VIDEO (Transcription)
    # ============================================

    def _extract_audio(self, path: Path) -> str:
        """Extract text from audio using Whisper"""
        try:
            import whisper

            print(f"🎵 Transcribing audio: {path.name}...")
            print("⏳ This may take a few minutes for long files...")

            # Load model
            model = whisper.load_model("base")

            # Transcribe
            result = model.transcribe(str(path))

            transcription_text = result["text"]

            print(f"📝 TRANSCRIBED: {len(transcription_text)} characters")

            # Return JUST the text, no header
            return transcription_text

        except Exception as e:
            print(f"❌ Audio transcription failed: {e}")
            return f"[Audio: {path.name} - Transcription failed: {e}]"

    def _extract_video(self, path: Path) -> str:
        """Extract audio from video and transcribe"""
        try:
            from pydub import AudioSegment
            import whisper
            import tempfile
            import time
            import os

            print(f"🎬 Processing video: {path.name}...")

            # Create temp file
            temp_audio = tempfile.NamedTemporaryFile(suffix=".wav", delete=False)
            temp_path = temp_audio.name
            temp_audio.close()  # Close immediately so pydub can write to it

            try:
                # Extract audio
                print("  Extracting audio track...")
                video = AudioSegment.from_file(str(path))
                video.export(temp_path, format="wav")

                # Give the system time to release the file
                time.sleep(0.5)

                # Transcribe
                print("  Transcribing audio...")
                model = whisper.load_model("base")
                result = model.transcribe(temp_path)

                transcription = (
                    f"=== Transcription of {path.name} ===\n\n{result['text']}"
                )

                transcription = result["text"]  # Just the text
                return transcription  # No header

            finally:
                # Clean up - try multiple times if locked
                for i in range(5):
                    try:
                        if os.path.exists(temp_path):
                            os.unlink(temp_path)
                        break
                    except PermissionError:
                        if i < 4:
                            time.sleep(0.5)
                        else:
                            print(f"⚠️ Could not delete temp file: {temp_path}")

        except Exception as e:
            print(f"❌ Video transcription failed: {e}")
            return f"[Video: {path.name} - Transcription failed: {e}]"

    # ============================================
    # ARCHIVES
    # ============================================

    def _extract_zip(self, path: Path) -> str:
        """Extract and process all files in ZIP"""
        import zipfile
        import tempfile
        import shutil

        print(f"📦 Extracting ZIP: {path.name}...")

        text = f"=== Contents of {path.name} ===\n\n"

        with tempfile.TemporaryDirectory() as temp_dir:
            # Extract ZIP
            with zipfile.ZipFile(path, "r") as zip_ref:
                zip_ref.extractall(temp_dir)

            # Process each file
            for file_path in Path(temp_dir).rglob("*"):
                if file_path.is_file():
                    try:
                        # Recursively process each file
                        file_text = self.extract_text(str(file_path))
                        text += f"\n\n--- From: {file_path.name} ---\n\n"
                        text += "\n\n".join(file_text)
                    except Exception as e:
                        print(f"⚠️ Skipped {file_path.name}: {e}")
                        text += f"\n\n[Skipped: {file_path.name} - {e}]\n\n"

        return text

    def _extract_7z(self, path: Path) -> str:
        """Extract and process all files in 7Z"""
        import py7zr
        import tempfile

        print(f"📦 Extracting 7Z: {path.name}...")

        text = f"=== Contents of {path.name} ===\n\n"

        with tempfile.TemporaryDirectory() as temp_dir:
            # Extract 7Z
            with py7zr.SevenZipFile(path, "r") as archive:
                archive.extractall(temp_dir)

            # Process each file
            for file_path in Path(temp_dir).rglob("*"):
                if file_path.is_file():
                    try:
                        file_text = self.extract_text(str(file_path))
                        text += f"\n\n--- From: {file_path.name} ---\n\n"
                        text += "\n\n".join(file_text)
                    except Exception as e:
                        print(f"⚠️ Skipped {file_path.name}: {e}")
                        text += f"\n\n[Skipped: {file_path.name} - {e}]\n\n"

        return text

    # ============================================
    # STRUCTURED DATA
    # ============================================

    def _extract_html(self, path: Path) -> str:
        """Extract text from HTML"""
        from bs4 import BeautifulSoup

        with open(path, "r", encoding="utf-8") as f:
            soup = BeautifulSoup(f.read(), "html.parser")

            # Remove scripts and styles
            for script in soup(["script", "style"]):
                script.decompose()

            # Get text
            text = soup.get_text()

            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = "\n".join(chunk for chunk in chunks if chunk)

            return text

    def _extract_json(self, path: Path) -> str:
        """Extract text from JSON"""
        import json

        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)

        # Pretty print JSON as text
        return json.dumps(data, indent=2)

    def _extract_xml(self, path: Path) -> str:
        """Extract text from XML"""
        from bs4 import BeautifulSoup

        with open(path, "r", encoding="utf-8") as f:
            soup = BeautifulSoup(f.read(), "xml")
            return soup.get_text()

    # ============================================
    # CHUNKING
    # ============================================

    def _chunk_text(self, text: str) -> List[str]:
        """
        Split text into overlapping chunks

        Args:
            text: Full text

        Returns:
            List of chunks
        """
        words = text.split()
        chunks = []

        for i in range(0, len(words), self.chunk_size - self.chunk_overlap):
            chunk = " ".join(words[i : i + self.chunk_size])
            if chunk.strip():
                chunks.append(chunk)

        return chunks
